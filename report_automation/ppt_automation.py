
from pptx import Presentation
from pptx.util import Inches

prs = Presentation() # 파워포인트 객체

for i in range(0, 11): # 총 11개의 슬라이드 추가.
    title_slide_layout = prs.slide_layouts[i]
    slide = prs.slides.add_slide(title_slide_layout)

prs.save('add all slides.pptx')



